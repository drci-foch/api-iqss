"""
Générateur PowerPoint pour les indicateurs de Lettres de Liaison
Version 5.0 - Formatage proche du modèle Foch
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from datetime import datetime
from typing import Dict, List, Optional
import pandas as pd

# --------------------------------------------------------------------
#  CONSTANTES GLOBALES
# --------------------------------------------------------------------

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

# Couleurs Hôpital Foch (palette officielle)
FOCH_BLUE = RGBColor(0, 82, 147)  # #005293
FOCH_GREEN = RGBColor(106, 168, 79)  # #6AA84F
FOCH_LIGHT_BLUE = RGBColor(155, 194, 230)  # #9BC2E6
FOCH_DARK_BLUE = RGBColor(0, 51, 102)  # #003366
FOCH_GRAY = RGBColor(89, 89, 89)  # #595959

# Couleurs indicateurs
COLOR_GREEN = RGBColor(146, 208, 80)
COLOR_YELLOW = RGBColor(255, 192, 0)
COLOR_ORANGE = RGBColor(255, 127, 39)
COLOR_RED = RGBColor(255, 0, 0)
COLOR_GRAY = RGBColor(217, 217, 217)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)


# --------------------------------------------------------------------
#  FONCTIONS UTILITAIRES DE STYLE
# --------------------------------------------------------------------


def set_cell_border(cell, border_color=RGBColor(0, 0, 0), border_width=Pt(0.5)):
    """Ajouter des bordures à une cellule"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Créer les bordures
    for border_name in ["left", "right", "top", "bottom"]:
        border = OxmlElement(f"a:{border_name}")
        border.set("w", str(int(border_width)))
        border.set("cap", "flat")
        border.set("cmpd", "sng")
        border.set("algn", "ctr")

        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")

        solidFill.append(srgbClr)
        border.append(solidFill)

        tcPr.append(border)


def set_cell_style(
    cell,
    text,
    font_size=9,
    bold=False,
    fill_color=None,
    text_color=None,
    align_center=True,
    border=True,
):
    """Appliquer un style uniforme à une cellule avec bordures"""
    cell.text = str(text)

    # Texte
    for paragraph in cell.text_frame.paragraphs:
        if align_center:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        else:
            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0

        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = "Calibri"
            if text_color:
                run.font.color.rgb = text_color
            else:
                run.font.color.rgb = COLOR_BLACK

    # Alignement vertical (milieu)
    cell.vertical_anchor = 1  # Middle

    # Remplissage
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

    # Marges réduites
    cell.text_frame.margin_left = Pt(3)
    cell.text_frame.margin_right = Pt(3)
    cell.text_frame.margin_top = Pt(2)
    cell.text_frame.margin_bottom = Pt(2)
    cell.text_frame.word_wrap = True

    # Bordures
    if border:
        set_cell_border(cell, FOCH_GRAY, Pt(0.75))


def get_color_by_threshold(value: float, excellent=95, good=85, medium=70) -> RGBColor:
    """Obtenir la couleur selon les seuils"""
    if pd.isna(value):
        return COLOR_GRAY
    if value >= excellent:
        return COLOR_GREEN
    elif value >= good:
        return COLOR_YELLOW
    elif value >= medium:
        return COLOR_ORANGE
    else:
        return COLOR_RED


def add_bullet_paragraph(
    text_frame, text, level=0, font_size=14, bold=False, color=FOCH_DARK_BLUE
):
    """Ajouter un paragraphe avec puce"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_before = Pt(3)
    p.space_after = Pt(3)
    p.line_spacing = 1.15
    return p


def add_footer(
    slide,
    page_number: Optional[int] = None,
    footer_text: str = "Indicateur prioritaire délai de validation des lettres de liaison",
):
    """
    Ajoute un pied de page homogène sur toutes les diapos :
    - texte à gauche
    - numéro de page à droite
    """
    # Texte à gauche
    left_box = slide.shapes.add_textbox(
        Inches(0.4), SLIDE_HEIGHT - Inches(0.45), Inches(7), Inches(0.35)
    )
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(8)
    p.font.color.rgb = FOCH_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    # Numéro de page à droite
    if page_number is not None:
        page_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(0.7),
            SLIDE_HEIGHT - Inches(0.45),
            Inches(0.5),
            Inches(0.35),
        )
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(page_number)
        p.font.size = Pt(8)
        p.font.color.rgb = FOCH_GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.RIGHT


def add_logo(slide, logo_path: Optional[str], height_inch: float = 0.8):
    """
    Ajoute le logo Hôpital Foch en haut à droite si un chemin est fourni.
    """
    if not logo_path:
        return
    try:
        slide.shapes.add_picture(
            logo_path,
            SLIDE_WIDTH - Inches(2.6),
            Inches(0.4),
            height=Inches(height_inch),
        )
    except FileNotFoundError:
        # On ne casse pas la génération si le logo n'est pas trouvé
        pass


# --------------------------------------------------------------------
#  SLIDES
# --------------------------------------------------------------------


def create_slide_1_title(
    prs: Presentation, period: str, logo_path: Optional[str] = None
) -> None:
    """Slide 1 : Page de titre proche de la maquette PDF"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond blanc par défaut (style corporate)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_WHITE

    # Grande "goutte" colorée à gauche (approximation de la forme verte/bleue)
    shape_green = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, -Inches(3.5), Inches(0.4), Inches(8), Inches(8)
    )
    shape_green.fill.solid()
    shape_green.fill.fore_color.rgb = FOCH_GREEN
    shape_green.line.fill.background()

    shape_blue = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, -Inches(3.0), Inches(0.0), Inches(8), Inches(8)
    )
    shape_blue.fill.solid()
    shape_blue.fill.fore_color.rgb = FOCH_BLUE
    shape_blue.line.fill.background()

    # Logo en haut à droite (si fourni)
    add_logo(slide, logo_path, height_inch=0.9)

    # Bloc de titre à droite
    title_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.4), Inches(4.3), Inches(2.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = (
        "INDICATEURS PRIORITAIRES:\n"
        "délai de validation et de diffusion (envoi)\n"
        "des lettres de liaison (LL)\n"
        "des séjours\n"
        "> à 24 h (1 nuit et plus)"
    )
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.1

    # Sous-titre période
    p = tf.add_paragraph()
    p.text = "Améliorons ensemble nos résultats"
    p.font.size = Pt(12)
    p.font.color.rgb = FOCH_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(10)

    p = tf.add_paragraph()
    p.text = f"Résultats du {period}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    # Pied de page (page 1)
    add_footer(slide, page_number=1)


def create_slide_2_methodology_part1(
    prs: Presentation, logo_path: Optional[str] = None
) -> None:
    """Slide 2 : Méthodologie partie 1 - typologie des séjours"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Bandeau titre
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.7)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FOCH_BLUE
    header_shape.line.fill.background()

    # Titre bandeau
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.12), SLIDE_WIDTH - Inches(0.8), Inches(0.45)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Indicateur prioritaire : délai de validation des lettres de liaison"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Logo discret en haut à droite (optionnel)
    add_logo(slide, logo_path, height_inch=0.6)

    # Sous-titre
    subtitle_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4),
        Inches(0.85),
        SLIDE_WIDTH - Inches(0.8),
        Inches(0.35),
    )
    subtitle_shape.fill.solid()
    subtitle_shape.fill.fore_color.rgb = FOCH_LIGHT_BLUE
    subtitle_shape.line.fill.background()

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), SLIDE_WIDTH - Inches(1), Inches(0.3)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Typologie des séjours"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = FOCH_DARK_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    # Contenu
    content_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.35), SLIDE_WIDTH - Inches(1.2), Inches(4.0)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = (
        "❖ Le Décret n° 2016-995 du 20 juillet 2016 relatif aux lettres de liaison "
        "(NOR : AFSH1612283D) précise que lors de la sortie de l'établissement de santé, "
        "une lettre de liaison (LL), rédigée par le médecin de l'établissement qui l'a pris en charge, "
        "est remise au patient et transmise le même jour au médecin traitant."
    )
    p.font.size = Pt(11)
    p.font.color.rgb = FOCH_DARK_BLUE
    p.font.name = "Calibri"
    p.space_after = Pt(8)
    p.line_spacing = 1.15

    add_bullet_paragraph(
        tf,
        '❖ Le code de santé publique demande une LL à la sortie de toute "admission" '
        "(en opposition aux consultations), HDJ comprises",
        level=0,
        font_size=11,
    )

    p = add_bullet_paragraph(
        tf,
        "❖ Séjours pris en compte pour l'indicateur « séjours de 1 nuit et plus » :",
        level=0,
        font_size=11,
        bold=True,
    )

    add_bullet_paragraph(
        tf,
        "➢ Les séjours suivants sont exclus :",
        level=1,
        font_size=10,
        bold=True,
    )

    exclusions = [
        "▪ Patients décédés (séjours non soumis aux LL)",
        "▪ Chirurgie ambulatoire et Hôpitaux de jour",
        "▪ Anesthésie, ophtalmologie, radiologie, ORL 392A",
    ]
    for excl in exclusions:
        add_bullet_paragraph(tf, excl, level=2, font_size=10)

    # Pied de page (page 2)
    add_footer(slide, page_number=2)


def create_slide_3_methodology_part2(
    prs: Presentation, logo_path: Optional[str] = None
) -> None:
    """Slide 3 : Méthodologie partie 2 - principe des indicateurs de diffusion"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Bandeau titre
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.7)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FOCH_BLUE
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.12), SLIDE_WIDTH - Inches(0.8), Inches(0.45)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Indicateur prioritaire : délai de validation des lettres de liaison"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    add_logo(slide, logo_path, height_inch=0.6)

    # Sous-titre
    subtitle_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4),
        Inches(0.85),
        SLIDE_WIDTH - Inches(0.8),
        Inches(0.35),
    )
    subtitle_shape.fill.solid()
    subtitle_shape.fill.fore_color.rgb = FOCH_LIGHT_BLUE
    subtitle_shape.line.fill.background()

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.9), SLIDE_WIDTH - Inches(1), Inches(0.3)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Principe des indicateurs de diffusion (envois)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = FOCH_DARK_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    # Contenu
    content_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.35), SLIDE_WIDTH - Inches(1.2), Inches(3.2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "➢ Seuls les séjours avec lettre de liaison validée par le médecin sont pris en compte."
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = FOCH_DARK_BLUE
    p.font.name = "Calibri"
    p.space_after = Pt(8)

    add_bullet_paragraph(tf, "➢ En excluant :", level=0, font_size=11, bold=True)

    exclusions = [
        "▸ Les LL validées les samedis, dimanches et jours fériés "
        "(jours d'absence des secrétaires)",
        "▸ Les LL avec plusieurs versions, dont la dernière version est validée à partir de J+1 "
        "après la sortie (date de diffusion des versions antérieures non sauvegardées)",
    ]
    for excl in exclusions:
        add_bullet_paragraph(tf, excl, level=1, font_size=10)

    # Encadré note
    note_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(4.7),
        SLIDE_WIDTH - Inches(1.6),
        Inches(0.9),
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RGBColor(220, 230, 241)
    note_box.line.color.rgb = FOCH_BLUE
    note_box.line.width = Pt(2)

    tf = note_box.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ℹ️  NOTE IMPORTANTE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(3)

    p = tf.add_paragraph()
    p.text = (
        "L'indicateur de diffusion mesure le respect du décret sur l'envoi le jour même de la sortie, "
        "en tenant compte des contraintes organisationnelles (week-ends et jours fériés)."
    )
    p.font.size = Pt(10)
    p.font.color.rgb = FOCH_DARK_BLUE
    p.font.name = "Calibri"
    p.line_spacing = 1.15

    # Pied de page (page 3)
    add_footer(slide, page_number=3)


def create_slide_4_summary(
    prs: Presentation, stats: Dict, period: str, logo_path: Optional[str] = None
) -> None:
    """Slide 4 : Résumé global"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.6)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FOCH_BLUE
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.1), SLIDE_WIDTH - Inches(1), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"RÉSUMÉ GLOBAL - {period}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    add_logo(slide, logo_path, height_inch=0.6)

    # Tableau centré
    rows = 4
    cols = 4
    table = slide.shapes.add_table(
        rows, cols, Inches(1.5), Inches(1.2), SLIDE_WIDTH - Inches(3), Inches(3.0)
    ).table

    # Largeurs de colonnes
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(1.4)
    table.columns[3].width = Inches(0.8)

    headers = ["Indicateur", "Valeur", "Objectif", "Statut"]
    for col_idx, header in enumerate(headers):
        set_cell_style(
            table.cell(0, col_idx), header, 11, True, FOCH_DARK_BLUE, COLOR_WHITE
        )

    # Ligne 1
    set_cell_style(
        table.cell(1, 0), "Nombre total de séjours", 10, True, align_center=False
    )
    set_cell_style(
        table.cell(1, 1), f"{stats['total_sejours']:,}".replace(",", " "), 10
    )
    set_cell_style(table.cell(1, 2), "-", 10)
    set_cell_style(table.cell(1, 3), "📊", 12)

    # Ligne 2
    taux_val = stats["taux_validation"]
    color_val = get_color_by_threshold(taux_val, 95, 85, 70)
    statut_val = "✅" if taux_val >= 95 else "⚠️" if taux_val >= 85 else "❌"
    set_cell_style(table.cell(2, 0), "Taux de validation", 10, True, align_center=False)
    set_cell_style(
        table.cell(2, 1), f"{taux_val:.1f}%", 10, False, color_val, COLOR_BLACK
    )
    set_cell_style(table.cell(2, 2), "≥ 95%", 10)
    set_cell_style(table.cell(2, 3), statut_val, 14)

    # Ligne 3
    taux_j0 = stats["taux_validation_j0"]
    color_j0 = get_color_by_threshold(taux_j0, 90, 80, 70)
    statut_j0 = "✅" if taux_j0 >= 90 else "⚠️" if taux_j0 >= 80 else "❌"
    set_cell_style(table.cell(3, 0), "Taux validation J0", 10, True, align_center=False)
    set_cell_style(
        table.cell(3, 1), f"{taux_j0:.1f}%", 10, False, color_j0, COLOR_BLACK
    )
    set_cell_style(table.cell(3, 2), "≥ 90%", 10)
    set_cell_style(table.cell(3, 3), statut_j0, 14)

    # Légende
    legend_box = slide.shapes.add_textbox(
        Inches(0.5), SLIDE_HEIGHT - Inches(1.1), SLIDE_WIDTH - Inches(1), Inches(0.5)
    )
    tf = legend_box.text_frame
    p = tf.paragraphs[0]
    p.text = (
        "Légende : 🟢 Vert = Excellent (objectif atteint)  |  🟡 Jaune = Bon  |  "
        "🟠 Orange = Moyen  |  🔴 Rouge = À améliorer"
    )
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = FOCH_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Pied de page (page 4)
    add_footer(slide, page_number=4)


def create_slide_5_validation_table(
    prs: Presentation, stats: Dict, period: str, logo_path: Optional[str] = None
) -> None:
    """Slide 5 : Tableau détaillé par spécialité"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Bandeau titre
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.55)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FOCH_BLUE
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.25), Inches(0.08), SLIDE_WIDTH - Inches(0.5), Inches(0.4)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Taux de validation le jour de la sortie et de diffusion des lettres de liaison (LL)\n"
        f"SÉJOURS > 24 H - Résultats du {period}"
    )
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.1

    add_logo(slide, logo_path, height_inch=0.55)

    # Tableau
    specialites = stats["par_specialite"]
    rows = len(specialites) + 2  # + en-tête + total
    cols = 9

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.1),
        Inches(0.6),
        SLIDE_WIDTH - Inches(0.2),
        SLIDE_HEIGHT - Inches(1.4),
    ).table

    headers = [
        "SPÉCIALITÉS",
        "Nb\ntotal",
        "LL\nvalid.",
        "%\nval.",
        "%\nJ0",
        "Délai\nval. (j)",
        "LL\ndiff.",
        "%\ndiff.",
        "Délai\ndiff. (j)",
    ]

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        set_cell_style(
            cell, header, 7, True, FOCH_LIGHT_BLUE, FOCH_DARK_BLUE, True, True
        )

    # Largeurs
    col_widths = [
        Inches(1.6),
        Inches(0.6),
        Inches(0.6),
        Inches(0.7),
        Inches(0.7),
        Inches(0.7),
        Inches(0.7),
        Inches(0.7),
        Inches(0.7),
    ]
    for col_idx, width in enumerate(col_widths):
        table.columns[col_idx].width = width

    # Lignes de données
    for row_idx, spe in enumerate(specialites, start=1):
        row_fill = COLOR_WHITE if row_idx % 2 == 1 else RGBColor(242, 242, 242)

        # Spécialité
        set_cell_style(
            table.cell(row_idx, 0),
            spe["specialite"],
            7,
            True,
            row_fill,
            FOCH_DARK_BLUE,
            False,
            True,
        )

        set_cell_style(table.cell(row_idx, 1), spe["nb_total"], 7, fill_color=row_fill)
        set_cell_style(
            table.cell(row_idx, 2), spe["nb_valides"], 7, fill_color=row_fill
        )

        taux = spe["taux_validation"]
        color_val = get_color_by_threshold(taux, 95, 85, 70)
        set_cell_style(
            table.cell(row_idx, 3),
            f"{taux:.1f}%",
            7,
            False,
            color_val,
            COLOR_BLACK,
        )

        taux_j0 = spe["taux_validation_j0"]
        color_j0 = get_color_by_threshold(taux_j0, 90, 80, 70)
        set_cell_style(
            table.cell(row_idx, 4),
            f"{taux_j0:.1f}%",
            7,
            False,
            color_j0,
            COLOR_BLACK,
        )

        delai_val = spe.get("delai_moyen", 0)
        set_cell_style(
            table.cell(row_idx, 5),
            f"{delai_val:.1f}",
            7,
            fill_color=row_fill,
        )

        nb_diff = spe.get("nb_diffuses", spe["nb_valides"])
        set_cell_style(
            table.cell(row_idx, 6),
            nb_diff,
            7,
            fill_color=row_fill,
        )

        pct_diff = spe.get("pct_diffuses", 100)
        color_diff = get_color_by_threshold(pct_diff, 90, 75, 60)
        set_cell_style(
            table.cell(row_idx, 7),
            f"{pct_diff:.1f}%",
            7,
            False,
            color_diff,
            COLOR_BLACK,
        )

        delai_diff = spe.get("delai_diffusion", 0)
        set_cell_style(
            table.cell(row_idx, 8),
            f"{delai_diff:.1f}",
            7,
            fill_color=row_fill,
        )

    # Ligne TOTAL FOCH
    total_row = rows - 1
    for col_idx in range(cols):
        table.cell(total_row, col_idx).fill.solid()
        table.cell(total_row, col_idx).fill.fore_color.rgb = FOCH_DARK_BLUE

    set_cell_style(
        table.cell(total_row, 0),
        "TOTAL FOCH",
        8,
        True,
        FOCH_DARK_BLUE,
        COLOR_WHITE,
        False,
        True,
    )

    set_cell_style(
        table.cell(total_row, 1),
        f"{stats['total_sejours']:,}".replace(",", " "),
        8,
        True,
        FOCH_DARK_BLUE,
        COLOR_WHITE,
    )

    set_cell_style(
        table.cell(total_row, 2),
        f"{stats['sejours_valides']:,}".replace(",", " "),
        8,
        True,
        FOCH_DARK_BLUE,
        COLOR_WHITE,
    )

    taux_global = stats["taux_validation"]
    color_global = get_color_by_threshold(taux_global, 95, 85, 70)
    set_cell_style(
        table.cell(total_row, 3),
        f"{taux_global:.1f}%",
        8,
        True,
        color_global,
        COLOR_BLACK,
    )

    taux_j0_global = stats["taux_validation_j0"]
    color_j0_global = get_color_by_threshold(taux_j0_global, 90, 80, 70)
    set_cell_style(
        table.cell(total_row, 4),
        f"{taux_j0_global:.1f}%",
        8,
        True,
        color_j0_global,
        COLOR_BLACK,
    )

    delai_global = stats.get("delai_moyen_validation", 0)
    set_cell_style(
        table.cell(total_row, 5),
        f"{delai_global:.1f}",
        8,
        True,
        FOCH_DARK_BLUE,
        COLOR_WHITE,
    )

    total_diff = stats.get("total_diffuses", stats["sejours_valides"])
    set_cell_style(
        table.cell(total_row, 6),
        f"{total_diff:,}".replace(",", " "),
        8,
        True,
        FOCH_DARK_BLUE,
        COLOR_WHITE,
    )

    pct_diff_global = stats.get("pct_diffuses", 100)
    color_diff_global = get_color_by_threshold(pct_diff_global, 90, 75, 60)
    set_cell_style(
        table.cell(total_row, 7),
        f"{pct_diff_global:.1f}%",
        8,
        True,
        color_diff_global,
        COLOR_BLACK,
    )

    delai_diff_global = stats.get("delai_moyen_diffusion", 0)
    set_cell_style(
        table.cell(total_row, 8),
        f"{delai_diff_global:.1f}",
        8,
        True,
        FOCH_DARK_BLUE,
        COLOR_WHITE,
    )

    # Pied de page (page 5)
    add_footer(slide, page_number=5)


def create_slide_6_instructions(
    prs: Presentation, logo_path: Optional[str] = None
) -> None:
    """Slide 6 : Rappel procédure / objectifs"""

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.6)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = FOCH_BLUE
    header_shape.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.1), SLIDE_WIDTH - Inches(1), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RAPPEL : VALIDATION DES LETTRES DE LIAISON"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    add_logo(slide, logo_path, height_inch=0.6)

    # Colonne gauche : procédure
    left_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1), Inches(4.3), Inches(3.8)
    )
    tf = left_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "✓ Procédure de validation :"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.font.name = "Calibri"
    p.space_after = Pt(6)

    instructions = [
        "1. Valider la lettre LE JOUR DE LA SORTIE",
        "2. Vérifier les informations complètes",
        "3. Cliquer sur « Valider » dans EASILY",
        "4. Diffusion automatique aux destinataires",
    ]
    for instruction in instructions:
        add_bullet_paragraph(tf, instruction, level=0, font_size=11)

    # Colonne droite : objectifs
    right_box = slide.shapes.add_textbox(
        Inches(5.1), Inches(1), Inches(4.3), Inches(3.8)
    )
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "⚠️  OBJECTIFS À ATTEINDRE :"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.font.name = "Calibri"
    p.space_after = Pt(6)

    objectifs = [
        "• Taux validation : ≥ 95%",
        "• Taux validation J0 : ≥ 90%",
        "• Respect du décret 2016-995",
    ]
    for obj in objectifs:
        add_bullet_paragraph(tf, obj, level=0, font_size=11, bold=True)

    # Encadré contact
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5),
        SLIDE_HEIGHT - Inches(1.45),
        Inches(7),
        Inches(0.5),
    )
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = FOCH_GREEN
    contact_box.line.color.rgb = FOCH_DARK_BLUE
    contact_box.line.width = Pt(2)

    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = (
        "📞 Contact : Gaëlle BURDY – Direction Qualité – DECT 2105 – "
        "gaelle.burdy@hopital-foch.com"
    )
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Pied de page (page 6)
    add_footer(slide, page_number=6)


# --------------------------------------------------------------------
#  GENERATION DU PPT
# --------------------------------------------------------------------


def generate_powerpoint(
    stats_validation: Dict,
    stats_diffusion: Dict,  # gardé pour compatibilité si tu l'utilises ailleurs
    output_path: str,
    period: str,
    logo_path: Optional[str] = None,
) -> None:
    """Générer le PowerPoint avec un formatage proche de la maquette Foch"""

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    create_slide_1_title(prs, period, logo_path=logo_path)
    create_slide_2_methodology_part1(prs, logo_path=logo_path)
    create_slide_3_methodology_part2(prs, logo_path=logo_path)
    create_slide_4_summary(prs, stats_validation, period, logo_path=logo_path)
    create_slide_5_validation_table(prs, stats_validation, period, logo_path=logo_path)
    create_slide_6_instructions(prs, logo_path=logo_path)

    prs.save(output_path)
    print(f"✅ PowerPoint généré : {output_path}")
    print(f"   {len(prs.slides)} slides | Formatage harmonisé")


if __name__ == "__main__":
    # Exemple de test
    test_stats = {
        "total_sejours": 1769,
        "sejours_valides": 1603,
        "taux_validation": 90.6,
        "taux_validation_j0": 70.7,
        "delai_moyen_validation": 0.8,
        "total_diffuses": 1603,
        "pct_diffuses": 100.0,
        "delai_moyen_diffusion": 0.8,
        "par_specialite": [
            {
                "specialite": "VASCULAIRE",
                "nb_total": 128,
                "nb_valides": 117,
                "taux_validation": 91.4,
                "taux_validation_j0": 72.6,
                "delai_moyen": 0.8,
                "nb_diffuses": 117,
                "pct_diffuses": 100.0,
                "delai_diffusion": 0.8,
            },
            {
                "specialite": "NEUROCHIRURGIE",
                "nb_total": 145,
                "nb_valides": 140,
                "taux_validation": 96.5,
                "taux_validation_j0": 85.0,
                "delai_moyen": 0.5,
                "nb_diffuses": 140,
                "pct_diffuses": 100.0,
                "delai_diffusion": 0.5,
            },
            {
                "specialite": "CARDIOLOGIE",
                "nb_total": 197,
                "nb_valides": 180,
                "taux_validation": 91.4,
                "taux_validation_j0": 75.5,
                "delai_moyen": 0.7,
                "nb_diffuses": 180,
                "pct_diffuses": 100.0,
                "delai_diffusion": 0.7,
            },
        ],
    }

    generate_powerpoint(
        test_stats,
        test_stats,
        "test_output_v5.pptx",
        "01/01 au 31/07/2025 (TEST)",
        logo_path=None,  # mets ici le chemin du logo Foch si tu l'as, ex: "logo_foch.png"
    )
    print("\n✅ Test V5 terminé !")
