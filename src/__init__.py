"""
Générateur PowerPoint pour les indicateurs de Lettres de Liaison
Version 2.0 - Avec plusieurs slides et tableaux améliorés
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from datetime import datetime
from typing import Dict, List
import pandas as pd

# Couleurs du thème Hôpital Foch
FOCH_BLUE = RGBColor(0, 82, 147)  # #005293
FOCH_GREEN = RGBColor(106, 168, 79)  # #6AA84F
FOCH_LIGHT_BLUE = RGBColor(155, 194, 230)  # #9BC2E6
FOCH_DARK_BLUE = RGBColor(0, 51, 102)  # Plus foncé pour contraste

# Couleurs pour les indicateurs
COLOR_GREEN = RGBColor(146, 208, 80)  # Vert : Excellent (≥95%)
COLOR_YELLOW = RGBColor(255, 192, 0)  # Jaune : Bon (≥85%)
COLOR_ORANGE = RGBColor(255, 127, 39)  # Orange : Moyen (≥70%)
COLOR_RED = RGBColor(255, 0, 0)  # Rouge : Faible (<70%)
COLOR_GRAY = RGBColor(200, 200, 200)  # Gris : Pas de données

# Blanc pour le texte
COLOR_WHITE = RGBColor(255, 255, 255)


def set_cell_style(
    cell, text, font_size=9, bold=False, fill_color=None, text_color=None
):
    """Appliquer un style uniforme à une cellule"""
    cell.text = str(text)

    # Texte
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = "Calibri"
            if text_color:
                run.font.color.rgb = text_color

    # Remplissage
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

    # Marges
    cell.text_frame.margin_left = Pt(2)
    cell.text_frame.margin_right = Pt(2)
    cell.text_frame.margin_top = Pt(2)
    cell.text_frame.margin_bottom = Pt(2)
    cell.text_frame.word_wrap = True


def get_color_by_threshold(value: float, excellent=95, good=85, medium=70) -> RGBColor:
    """Obtenir la couleur selon les seuils"""
    if pd.isna(value):
        return COLOR_GRAY
    if value >= excellent:
        return COLOR_GREEN
    elif value >= good:
        return COLOR_YELLOW
    elif value >= medium:
        return COLOR_ORANGE
    else:
        return COLOR_RED


def create_slide_1_title(prs: Presentation, period: str) -> None:
    """
    Slide 1 : Page de titre
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "INDICATEURS PRIORITAIRES\n\n"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    p = tf.add_paragraph()
    p.text = "Délai de validation et de diffusion (envoi)\ndes lettres de liaison (LL)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Mention séjours
    p = tf.add_paragraph()
    p.text = "\ndes séjours > à 24 h (1 nuit et plus)"
    p.font.size = Pt(20)
    p.font.color.rgb = FOCH_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Période
    period_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = period_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Améliorons ensemble nos résultats\nRésultats du {period}"
    p.font.size = Pt(18)
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER


def create_slide_2_summary(prs: Presentation, stats: Dict, period: str) -> None:
    """
    Slide 2 : Résumé général avec indicateurs clés
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"RÉSUMÉ GLOBAL - {period}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Tableau résumé
    rows = 4
    cols = 4
    table = slide.shapes.add_table(
        rows, cols, Inches(2), Inches(1.5), Inches(6), Inches(3)
    ).table

    # En-tête
    headers = ["Indicateur", "Valeur", "Objectif", "Statut"]
    for col_idx, header in enumerate(headers):
        set_cell_style(
            table.cell(0, col_idx), header, 11, True, FOCH_DARK_BLUE, COLOR_WHITE
        )

    # Ligne 1 : Nombre total de séjours
    set_cell_style(table.cell(1, 0), "Nombre total de séjours", 10, True)
    set_cell_style(table.cell(1, 1), f"{stats['total_sejours']:,}", 10)
    set_cell_style(table.cell(1, 2), "-", 10)
    set_cell_style(table.cell(1, 3), "📊", 10)

    # Ligne 2 : Taux de validation
    taux_val = stats["taux_validation"]
    color_val = get_color_by_threshold(taux_val, 95, 85, 70)
    statut_val = "✅" if taux_val >= 95 else "⚠️" if taux_val >= 85 else "❌"
    set_cell_style(table.cell(2, 0), "Taux de validation", 10, True)
    set_cell_style(table.cell(2, 1), f"{taux_val}%", 10, fill_color=color_val)
    set_cell_style(table.cell(2, 2), "≥ 95%", 10)
    set_cell_style(table.cell(2, 3), statut_val, 10)

    # Ligne 3 : Taux validation J0
    taux_j0 = stats["taux_validation_j0"]
    color_j0 = get_color_by_threshold(taux_j0, 90, 80, 70)
    statut_j0 = "✅" if taux_j0 >= 90 else "⚠️" if taux_j0 >= 80 else "❌"
    set_cell_style(table.cell(3, 0), "Taux validation J0", 10, True)
    set_cell_style(table.cell(3, 1), f"{taux_j0}%", 10, fill_color=color_j0)
    set_cell_style(table.cell(3, 2), "≥ 90%", 10)
    set_cell_style(table.cell(3, 3), statut_j0, 10)

    # Note en bas
    note_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.8))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Couleurs : Vert = Excellent | Jaune = Bon | Orange = Moyen | Rouge = À améliorer"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = FOCH_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER


def create_slide_3_validation_table(
    prs: Presentation, stats: Dict, period: str
) -> None:
    """
    Slide 3 : Tableau détaillé VALIDATION par spécialité
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.25), Inches(9.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"INDICATEURS DE VALIDATION PAR SPÉCIALITÉ - {period}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Tableau
    specialites = stats["par_specialite"]
    rows = len(specialites) + 2  # +1 en-tête, +1 total
    cols = 6

    table = slide.shapes.add_table(
        rows, cols, Inches(0.3), Inches(0.9), Inches(9.4), Inches(5.8)
    ).table

    # En-têtes
    headers = [
        "SPÉCIALITÉ",
        "Nb total\nséjours",
        "LL\nvalidées",
        "% validés",
        "% validés\nJ0",
        "Délai moyen\nvalidation (j)",
    ]
    for col_idx, header in enumerate(headers):
        set_cell_style(table.cell(0, col_idx), header, 9, True, FOCH_LIGHT_BLUE)

    # Données
    for row_idx, spe in enumerate(specialites, start=1):
        # Spécialité
        set_cell_style(table.cell(row_idx, 0), spe["specialite"], 9, True)

        # Nb total
        set_cell_style(table.cell(row_idx, 1), spe["nb_total"], 9)

        # LL validées
        set_cell_style(table.cell(row_idx, 2), spe["nb_valides"], 9)

        # % validés avec couleur
        taux = spe["taux_validation"]
        color = get_color_by_threshold(taux, 95, 85, 70)
        set_cell_style(table.cell(row_idx, 3), f"{taux}%", 9, False, color)

        # % J0 avec couleur
        taux_j0 = spe["taux_validation_j0"]
        color_j0 = get_color_by_threshold(taux_j0, 90, 80, 70)
        set_cell_style(table.cell(row_idx, 4), f"{taux_j0}%", 9, False, color_j0)

        # Délai moyen
        delai = spe.get("delai_moyen", 0)
        set_cell_style(table.cell(row_idx, 5), f"{delai:.1f}", 9)

    # Ligne TOTAL
    total_row = rows - 1
    set_cell_style(
        table.cell(total_row, 0), "TOTAL FOCH", 10, True, FOCH_DARK_BLUE, COLOR_WHITE
    )
    set_cell_style(table.cell(total_row, 1), stats["total_sejours"], 10, True)
    set_cell_style(table.cell(total_row, 2), stats["sejours_valides"], 10, True)

    taux_global = stats["taux_validation"]
    color_global = get_color_by_threshold(taux_global, 95, 85, 70)
    set_cell_style(table.cell(total_row, 3), f"{taux_global}%", 10, True, color_global)

    taux_j0_global = stats["taux_validation_j0"]
    color_j0_global = get_color_by_threshold(taux_j0_global, 90, 80, 70)
    set_cell_style(
        table.cell(total_row, 4), f"{taux_j0_global}%", 10, True, color_j0_global
    )

    delai_global = stats.get("delai_moyen_validation", 0)
    set_cell_style(table.cell(total_row, 5), f"{delai_global:.1f}", 10, True)


def create_slide_4_diffusion_table(
    prs: Presentation, stats: Dict, stats_diffusion: Dict, period: str
) -> None:
    """
    Slide 4 : Tableau détaillé DIFFUSION par spécialité
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.25), Inches(9.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"INDICATEURS DE DIFFUSION PAR SPÉCIALITÉ - {period}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Tableau
    specialites = stats["par_specialite"]
    rows = len(specialites) + 2
    cols = 5

    table = slide.shapes.add_table(
        rows, cols, Inches(1.5), Inches(0.9), Inches(7), Inches(5.8)
    ).table

    # En-têtes
    headers = [
        "SPÉCIALITÉ",
        "LL validées",
        "LL diffusées",
        "% diffusées /\nvalidées",
        "Délai diffusion /\nvalidation (j)",
    ]
    for col_idx, header in enumerate(headers):
        set_cell_style(table.cell(0, col_idx), header, 9, True, FOCH_LIGHT_BLUE)

    # Données
    for row_idx, spe in enumerate(specialites, start=1):
        # Spécialité
        set_cell_style(table.cell(row_idx, 0), spe["specialite"], 9, True)

        # LL validées
        set_cell_style(table.cell(row_idx, 1), spe["nb_valides"], 9)

        # LL diffusées
        nb_diff = spe.get("nb_diffuses", spe["nb_valides"])
        set_cell_style(table.cell(row_idx, 2), nb_diff, 9)

        # % diffusées avec couleur
        pct_diff = spe.get("pct_valides", 100)
        color_diff = get_color_by_threshold(pct_diff, 90, 75, 60)
        set_cell_style(table.cell(row_idx, 3), f"{pct_diff}%", 9, False, color_diff)

        # Délai diffusion
        delai_diff = spe.get("delai_moyen", 0)
        set_cell_style(table.cell(row_idx, 4), f"{delai_diff:.1f}", 9)

    # Ligne TOTAL
    total_row = rows - 1
    set_cell_style(
        table.cell(total_row, 0), "TOTAL FOCH", 10, True, FOCH_DARK_BLUE, COLOR_WHITE
    )
    set_cell_style(table.cell(total_row, 1), stats["sejours_valides"], 10, True)

    total_diff = stats_diffusion.get("total_diffuses", stats["sejours_valides"])
    set_cell_style(table.cell(total_row, 2), total_diff, 10, True)

    pct_diff_global = stats_diffusion.get("pct_diffuses", 100)
    color_diff_global = get_color_by_threshold(pct_diff_global, 90, 75, 60)
    set_cell_style(
        table.cell(total_row, 3), f"{pct_diff_global}%", 10, True, color_diff_global
    )

    delai_diff_global = stats_diffusion.get("delai_moyen_diffusion", 0)
    set_cell_style(table.cell(total_row, 4), f"{delai_diff_global:.1f}", 10, True)


def create_slide_5_instructions(prs: Presentation) -> None:
    """
    Slide 5 : Instructions pour la validation des LL
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RAPPEL : COMMENT VALIDER UNE LETTRE DE LIAISON ?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Instructions
    instructions = [
        "1. Valider la lettre de liaison LE JOUR DE LA SORTIE du patient",
        "2. Vérifier que toutes les informations sont complètes et exactes",
        "3. Cliquer sur le bouton 'Valider' dans le système",
        "4. La lettre sera automatiquement diffusée aux destinataires",
        "",
        "⚠️  IMPORTANT :",
        "• Objectif : ≥ 95% de lettres validées",
        "• Objectif : ≥ 90% de lettres validées le jour de la sortie (J0)",
        "• En cas de difficulté, contacter la Direction Qualité",
    ]

    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf = text_box.text_frame
    tf.word_wrap = True

    for instruction in instructions:
        p = tf.add_paragraph() if instruction != instructions[0] else tf.paragraphs[0]
        p.text = instruction
        p.font.size = Pt(16)
        p.font.color.rgb = FOCH_DARK_BLUE
        p.level = 0
        if instruction.startswith("⚠️"):
            p.font.bold = True
            p.font.color.rgb = FOCH_BLUE
        if instruction.startswith("•"):
            p.level = 1

    # Contact
    contact_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.8), Inches(8), Inches(0.8)
    )
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Contact : Gaëlle BURDY - Direction Qualité - DECT 2105\ngaelle.burdy@hopital-foch.com"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = FOCH_BLUE
    p.alignment = PP_ALIGN.CENTER


def generate_powerpoint(
    stats_validation: Dict, stats_diffusion: Dict, output_path: str, period: str
) -> None:
    """
    Générer le PowerPoint complet avec toutes les slides

    Args:
        stats_validation: Statistiques de validation
        stats_diffusion: Statistiques de diffusion
        output_path: Chemin du fichier de sortie
        period: Période du rapport (ex: "01/01 au 31/07/2025")
    """

    # Créer la présentation
    prs = Presentation()

    # Définir la taille (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1 : Titre
    create_slide_1_title(prs, period)

    # Slide 2 : Résumé global
    create_slide_2_summary(prs, stats_validation, period)

    # Slide 3 : Tableau validation détaillé
    create_slide_3_validation_table(prs, stats_validation, period)

    # Slide 4 : Tableau diffusion détaillé
    create_slide_4_diffusion_table(prs, stats_validation, stats_diffusion, period)

    # Slide 5 : Instructions
    create_slide_5_instructions(prs)

    # Sauvegarder
    prs.save(output_path)


if __name__ == "__main__":
    # Test avec données fictives
    test_stats = {
        "total_sejours": 1500,
        "sejours_valides": 1380,
        "taux_validation": 92.0,
        "taux_validation_j0": 75.5,
        "delai_moyen_validation": 0.8,
        "par_specialite": [
            {
                "specialite": "CARDIOLOGIE",
                "nb_total": 200,
                "nb_valides": 185,
                "taux_validation": 92.5,
                "taux_validation_j0": 78.0,
                "delai_moyen": 0.7,
                "nb_diffuses": 185,
                "pct_valides": 100.0,
            },
            {
                "specialite": "NEUROCHIRURGIE",
                "nb_total": 150,
                "nb_valides": 145,
                "taux_validation": 96.7,
                "taux_validation_j0": 85.0,
                "delai_moyen": 0.5,
                "nb_diffuses": 145,
                "pct_valides": 100.0,
            },
            {
                "specialite": "VASCULAIRE",
                "nb_total": 180,
                "nb_valides": 160,
                "taux_validation": 88.9,
                "taux_validation_j0": 70.0,
                "delai_moyen": 1.2,
                "nb_diffuses": 160,
                "pct_valides": 100.0,
            },
        ],
    }

    test_diffusion = {
        "total_diffuses": 1380,
        "pct_diffuses": 100.0,
        "delai_moyen_diffusion": 0.8,
    }

    generate_powerpoint(
        test_stats, test_diffusion, "test_output.pptx", "01/01 au 31/07/2025 (TEST)"
    )
